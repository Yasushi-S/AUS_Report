from flask import Flask, request, render_template, redirect, url_for, flash
from dotenv import load_dotenv
import pandas as pd
from datetime import datetime
from pptx import Presentation
from pptx.util import Pt
import os
import re
import sqlite3
import sys
import logging
from config import Config

load_dotenv()

try:
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')
except (AttributeError, ValueError):
    pass

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
LOG_PATH = os.path.join(BASE_DIR, 'app.log')

logging.basicConfig(
    filename=LOG_PATH,
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)

app = Flask(__name__)
app.secret_key = Config.SECRET_KEY

DB_PATH = os.path.join(BASE_DIR, Config.DB_PATH)
TEMPLATE_PATH = os.path.join(BASE_DIR, Config.TEMPLATE_PATH)
IMPLEMENTATION_TEMPLATE_PATH = os.path.join(BASE_DIR, Config.IMPLEMENTATION_TEMPLATE_PATH)
OUTPUT_DIR = os.path.join(BASE_DIR, Config.OUTPUT_DIR)
PPTX_OUTPUT_DIR = os.path.join(BASE_DIR, Config.PPTX_OUTPUT_DIR)
CHANGES_LOG = os.path.join(BASE_DIR, 'changes.log')

os.makedirs(OUTPUT_DIR, exist_ok=True)
os.makedirs(PPTX_OUTPUT_DIR, exist_ok=True)


def get_db():
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    return conn


def init_db():
    with get_db() as conn:
        conn.execute('''
            CREATE TABLE IF NOT EXISTS cases (
                id              INTEGER PRIMARY KEY AUTOINCREMENT,
                seq_no          INTEGER NOT NULL,
                age             INTEGER NOT NULL,
                address         TEXT    NOT NULL,
                gestational_age INTEGER NOT NULL,
                date            TEXT    NOT NULL,
                patient_id      TEXT    NOT NULL,
                note            TEXT    NOT NULL DEFAULT '',
                UNIQUE(patient_id, date)
            )
        ''')


def row_to_dict(row):
    return {
        'id': row['id'],
        '通算No': row['seq_no'],
        'Age': row['age'],
        'Address': row['address'],
        'Gestational Age': row['gestational_age'],
        'Date': row['date'],
        'ID': row['patient_id'],
        '特記事項': row['note'] or '',
    }


def get_kankatsu_options():
    with get_db() as conn:
        rows = conn.execute(
            'SELECT address, COUNT(*) AS cnt FROM cases GROUP BY address ORDER BY cnt DESC'
        ).fetchall()
    return [row['address'] for row in rows]


def get_next_seq_no():
    with get_db() as conn:
        row = conn.execute('SELECT MAX(seq_no) AS max_no FROM cases').fetchone()
    return (row['max_no'] or 0) + 1


def get_available_months():
    with get_db() as conn:
        rows = conn.execute(
            "SELECT DISTINCT substr(date, 1, 7) AS month FROM cases ORDER BY month DESC"
        ).fetchall()
    months = [row['month'] for row in rows]
    if not months:
        months = [datetime.now().strftime('%Y-%m')]
    return months


def get_monthly_data(month):
    with get_db() as conn:
        rows = conn.execute(
            '''SELECT * FROM cases
               WHERE date LIKE ?
               ORDER BY date, patient_id''',
            (month + '%',)
        ).fetchall()
    data = []
    for row in rows:
        d = row_to_dict(row)
        d['has_note'] = bool(d['特記事項'])
        data.append(d)
    return data


def get_case(record_id):
    with get_db() as conn:
        row = conn.execute(
            'SELECT * FROM cases WHERE id = ?',
            (record_id,)
        ).fetchone()
    if row is None:
        return None
    return row_to_dict(row)


def validate_form_data(form_data):
    errors = []
    required_fields = ['通算No', 'Date', 'Age', 'Address', 'Gestational Age', 'ID']
    for field in required_fields:
        if not form_data.get(field):
            errors.append(f'{field}は必須項目です')

    try:
        if form_data.get('Age'):
            age = int(form_data['Age'])
            if age < 0 or age > 100:
                errors.append('年齢は0-100の範囲で入力してください')

        if form_data.get('Gestational Age'):
            ga = int(form_data['Gestational Age'])
            if ga < 0 or ga > 45:
                errors.append('妊娠週数は0-45の範囲で入力してください')

        if form_data.get('ID'):
            patient_id = form_data['ID'].zfill(5)
            if not re.fullmatch(r'\d{5}', patient_id):
                errors.append('IDは5桁の数字で入力してください')

        if form_data.get('Date'):
            try:
                date = datetime.strptime(form_data['Date'], '%Y-%m-%d')
                if date > datetime.now():
                    errors.append('未来の日付は入力できません')
            except ValueError:
                errors.append('日付の形式が正しくありません')
    except ValueError:
        errors.append('入力値の形式が正しくありません')

    return errors


def duplicate_exists(patient_id, date, exclude_patient_id=None, exclude_date=None):
    with get_db() as conn:
        if exclude_patient_id is not None and exclude_date is not None:
            row = conn.execute(
                '''SELECT 1 FROM cases
                   WHERE patient_id = ? AND date = ?
                   AND NOT (patient_id = ? AND date = ?)''',
                (patient_id, date, exclude_patient_id, exclude_date)
            ).fetchone()
        else:
            row = conn.execute(
                'SELECT 1 FROM cases WHERE patient_id = ? AND date = ?',
                (patient_id, date)
            ).fetchone()
    return row is not None


def log_deletion(case):
    try:
        with open(CHANGES_LOG, 'a', encoding='utf-8') as f:
            f.write(f"\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')} - 削除\n")
            f.write(f"seq_no: {case['seq_no']}\n")
            f.write(f"patient_id: {case['patient_id']}\n")
            f.write(f"date: {case['date']}\n")
            f.write(f"age: {case['age']}\n")
            f.write(f"address: {case['address']}\n")
            f.write(f"gestational_age: {case['gestational_age']}\n")
            f.write(f"note: {case['note']}\n")
    except Exception as e:
        logging.error(f"削除ログ記録エラー: {str(e)}")


def convert_gestational_age(age):
    """レポート生成時の週数変換（PPT用）"""
    if 5 <= age <= 7:
        return '1'
    elif 8 <= age <= 11:
        return '2'
    return ''


def convert_gestational_age_for_excel(age):
    """Excel出力用の週数変換（現行アプリ互換）"""
    if 5 <= age <= 7:
        return '1'
    elif 8 <= age <= 11:
        return '2'
    elif 12 <= age <= 15:
        return '3'
    elif 16 <= age <= 19:
        return '4'
    elif 20 <= age <= 21:
        return '5'
    return ''


def parse_address(address):
    """住所文字列を (都道府県, 市/郡地名, '市'/'郡', 町/村名) に分解する"""
    pref = ''
    city = ''
    gunshi = ''
    town = ''
    m = re.match(r'(.+?県|.+?都|.+?府|.+?道)', address)
    if m:
        pref = m.group(1).replace('県', '').replace('都', '').replace('府', '').replace('道', '')
        rest = address[m.end():]
    else:
        rest = address
    m = re.match(r'(.+?市|.+?郡)', rest)
    if m:
        city_full = m.group(1)
        if '市' in city_full:
            city = city_full.replace('市', '')
            gunshi = '市'
        elif '郡' in city_full:
            city = city_full.replace('郡', '')
            gunshi = '郡'
        rest = rest[m.end():]
    m = re.match(r'(.+?町|.+?村)', rest)
    if m:
        town = m.group(1).replace('町', '').replace('村', '')
    return pref, city, gunshi, town


def build_replacements(row, case_num, reiwa_year, report_month):
    n = str(case_num)
    date_obj = pd.to_datetime(row['date'])
    ga_raw = int(row['gestational_age'])
    ga = convert_gestational_age(ga_raw)
    pref, city, gunshi, town = parse_address(row['address'])

    return {
        f'RY{n}': (str(reiwa_year), 10),
        f'MM{n}': (str(report_month), 10),
        f'N{n}': (str(row['number']), 12),
        f'A{n}': (str(row['age']), 12),
        f'AD{n}-1': (pref, 12),
        f'AD{n}-2': (city, 12),
        f'AD{n}-3': (town, 12),
        f'M{n}': (str(date_obj.month), 12),
        f'D{n}': (str(date_obj.day), 12),
        f'G{n}-1': ('〇' if ga == '1' else '', 14),
        f'G{n}-2': ('〇' if ga == '2' else '', 14),
        f'GUN{n}': ('〇' if gunshi == '郡' else '', 14),
        f'SHI{n}': ('〇' if gunshi == '市' else '', 14),
        f'MACHI{n}': ('〇' if town else '', 14),
    }


def apply_replacements(slide, replacements):
    sorted_keys = sorted(replacements.keys(), key=len, reverse=True)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for key in sorted_keys:
                if key in para.text:
                    value, font_pt = replacements[key]
                    new_text = para.text.replace(key, value)
                    para.clear()
                    run = para.add_run()
                    run.text = new_text
                    run.font.size = Pt(font_pt)
                    break


def generate_ppt_reports(data, reiwa_year, month):
    if not os.path.exists(TEMPLATE_PATH):
        raise FileNotFoundError(f"テンプレートファイルが見つかりません: {TEMPLATE_PATH}")
    if not os.path.exists(IMPLEMENTATION_TEMPLATE_PATH):
        raise FileNotFoundError(f"実施数報告テンプレートファイルが見つかりません: {IMPLEMENTATION_TEMPLATE_PATH}")

    prs = Presentation(TEMPLATE_PATH)
    row_index = 0

    for slide in prs.slides:
        for i in range(2):
            if row_index >= len(data):
                break
            case_num = i + 1
            row_data = data.iloc[row_index]
            replacements = build_replacements(row_data, case_num, reiwa_year, month)
            apply_replacements(slide, replacements)
            row_index += 1

    required_slides = (len(data) + 1) // 2
    while len(prs.slides) > required_slides:
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[-1])

    output_ppt = os.path.join(PPTX_OUTPUT_DIR, f"AUS報告書_R{reiwa_year}_{month}月.pptx")
    prs.save(output_ppt)

    current_date = datetime.now()
    last_number = str(data['number'].iloc[-1]) if not data.empty else 'N/A'

    prs2 = Presentation(IMPLEMENTATION_TEMPLATE_PATH)
    for slide in prs2.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                replacements = {
                    'G1': str(reiwa_year),
                    'H1': str(month),
                    'I1': str(current_date.month),
                    'J1': str(current_date.day),
                    'K1': last_number,
                }
                for key, value in replacements.items():
                    if key in paragraph.text:
                        paragraph.text = paragraph.text.replace(key, value)

    output_implementation_ppt = os.path.join(PPTX_OUTPUT_DIR, f"AUS実施数報告_R{reiwa_year}_{month}月.pptx")
    prs2.save(output_implementation_ppt)
    logging.info(f"PowerPointファイル生成成功: 令和{reiwa_year}年{month}月")


@app.route('/', methods=['GET'])
def index():
    try:
        months = get_available_months()
        selected_month = request.args.get('month') or months[0]
        monthly_data = get_monthly_data(selected_month)
        return render_template(
            'index.html',
            kankatsu_options=get_kankatsu_options(),
            next_no=get_next_seq_no(),
            selected_month=selected_month,
            months=months,
            monthly_data=monthly_data,
        )
    except Exception as e:
        logging.error(f"インデックスページ表示エラー: {str(e)}")
        flash('データの読み込み中にエラーが発生しました。', 'error')
        return render_template('error.html')


@app.route('/edit/<int:record_id>', methods=['GET'])
def edit(record_id):
    try:
        patient_data = get_case(record_id)
        if patient_data is None:
            flash('指定されたデータが見つかりません。', 'error')
            return redirect(url_for('index'))
        return render_template(
            'edit.html',
            kankatsu_options=get_kankatsu_options(),
            patient_data=patient_data,
            record_id=record_id,
            weeks_range=range(5, 22),
        )
    except Exception as e:
        logging.error(f"編集ページ表示エラー: {str(e)}")
        flash('データの読み込み中にエラーが発生しました。', 'error')
        return redirect(url_for('index'))


@app.route('/submit_edit/<int:record_id>', methods=['POST'])
def submit_edit(record_id):
    try:
        errors = validate_form_data(request.form)
        if errors:
            for error in errors:
                flash(error, 'error')
            return redirect(url_for('edit', record_id=record_id))

        seq_no = int(request.form['通算No'])
        new_date = request.form['Date']
        age = int(request.form['Age'])
        address = request.form['Address']
        gestational_age = int(request.form['Gestational Age'])
        new_patient_id = request.form['ID'].zfill(5)
        note = request.form.get('特記事項', '')

        existing = get_case(record_id)
        if existing and duplicate_exists(new_patient_id, new_date, existing['ID'], existing['Date']):
            flash('同じ日付で既に同じ患者IDのデータが存在します', 'error')
            return redirect(url_for('edit', record_id=record_id))

        with get_db() as conn:
            conn.execute(
                '''UPDATE cases
                   SET seq_no = ?, age = ?, address = ?, gestational_age = ?,
                       date = ?, patient_id = ?, note = ?
                   WHERE id = ?''',
                (seq_no, age, address, gestational_age, new_date, new_patient_id, note,
                 record_id)
            )

        logging.info(f"データ更新: ID {new_patient_id}, Date {new_date}")
        flash('データが正常に更新されました。', 'success')
    except Exception as e:
        logging.error(f"データ更新エラー: {str(e)}")
        flash('データの更新中にエラーが発生しました。', 'error')

    return redirect(url_for('index'))


@app.route('/submit', methods=['POST'])
def submit():
    try:
        errors = validate_form_data(request.form)
        if errors:
            for error in errors:
                flash(error, 'error')
            return redirect(url_for('index'))

        seq_no = int(request.form['通算No'])
        new_date = request.form['Date']
        age = int(request.form['Age'])
        address = request.form['Address']
        gestational_age = int(request.form['Gestational Age'])
        new_patient_id = request.form['ID'].zfill(5)
        note = request.form.get('特記事項', '')

        if duplicate_exists(new_patient_id, new_date):
            flash('同じ日付で既に同じ患者IDのデータが存在します', 'error')
            return redirect(url_for('index'))

        with get_db() as conn:
            conn.execute(
                '''INSERT INTO cases
                   (seq_no, age, address, gestational_age, date, patient_id, note)
                   VALUES (?, ?, ?, ?, ?, ?, ?)''',
                (seq_no, age, address, gestational_age, new_date, new_patient_id, note)
            )

        logging.info(f"新規データ登録: ID {new_patient_id}")
        flash('データが正常に登録されました。', 'success')
    except Exception as e:
        logging.error(f"データ登録エラー: {str(e)}")
        flash('データの登録中にエラーが発生しました。', 'error')

    return redirect(url_for('index'))


@app.route('/delete/<int:record_id>', methods=['POST'])
def delete(record_id):
    try:
        with get_db() as conn:
            row = conn.execute(
                'SELECT * FROM cases WHERE id = ?',
                (record_id,)
            ).fetchone()
            if row is None:
                flash('指定されたデータが見つかりません。', 'error')
                return redirect(url_for('index'))

            conn.execute(
                'DELETE FROM cases WHERE id = ?',
                (record_id,)
            )
            log_deletion(dict(row))

        logging.info(f"データ削除: ID {row['patient_id']}, Date {row['date']}")
        flash('データが正常に削除されました。', 'success')
    except Exception as e:
        logging.error(f"データ削除エラー: {str(e)}")
        flash('データの削除中にエラーが発生しました。', 'error')

    return redirect(url_for('index'))


@app.route('/generate_report', methods=['POST'])
def generate_report():
    try:
        year_month = request.form['report_month']
        year = int(year_month.split('-')[0])
        month = int(year_month.split('-')[1])
        reiwa_year = year - 2018

        with get_db() as conn:
            rows = conn.execute(
                '''SELECT * FROM cases
                   WHERE substr(date, 1, 4) = ? AND substr(date, 6, 2) = ?
                   ORDER BY date, patient_id''',
                (str(year), f'{month:02d}')
            ).fetchall()

        if not rows:
            flash(f'{year}年{month}月のデータが見つかりません。', 'error')
            return redirect(url_for('index'))

        records = []
        for i, row in enumerate(rows, start=1):
            pref, city, gunshi, town = parse_address(row['address'])
            records.append({
                '通算No': row['seq_no'],
                'number': i,
                'age': row['age'],
                'address': row['address'],
                'gestational_age': row['gestational_age'],
                'date': row['date'],
                'patient_id': row['patient_id'],
                'note': row['note'] or '',
                'Address1': pref,
                'Address2': city,
                '郡市': gunshi,
                'Address3': town,
            })

        filtered_df = pd.DataFrame(records)
        excel_df = filtered_df.copy()
        excel_df['Gestational Age'] = excel_df['gestational_age'].apply(convert_gestational_age_for_excel)
        excel_df = excel_df.rename(columns={
            'age': 'Age',
            'patient_id': 'ID',
            'date': 'Date',
            'note': '特記事項',
        })
        excel_df = excel_df[[
            '通算No', 'number', 'Age',
            'Address1', 'Address2', '郡市', 'Address3',
            'Gestational Age', 'Date', 'ID', '特記事項'
        ]]
        excel_df = excel_df.rename(columns={'number': 'Number'})

        output_excel = os.path.join(OUTPUT_DIR, f"R{reiwa_year}_{month}月.xlsx")
        with pd.ExcelWriter(output_excel) as writer:
            excel_df.to_excel(writer, index=False)

        generate_ppt_reports(filtered_df, reiwa_year, month)

        logging.info(f"レポート生成成功: {year}年{month}月")
        flash(f'令和{reiwa_year}年{month}月のレポートが正常に生成されました。', 'success')
    except Exception as e:
        logging.error(f"レポート生成エラー: {str(e)}")
        flash(f'エラーが発生しました: {str(e)}', 'error')

    return redirect(url_for('index'))


init_db()

if __name__ == '__main__':
    try:
        missing_files = []
        if not os.path.exists(TEMPLATE_PATH):
            missing_files.append(f"テンプレートファイル: {TEMPLATE_PATH}")
        if not os.path.exists(IMPLEMENTATION_TEMPLATE_PATH):
            missing_files.append(f"実施数報告テンプレートファイル: {IMPLEMENTATION_TEMPLATE_PATH}")

        if missing_files:
            print("警告: 以下のファイルが見つかりません:")
            for file in missing_files:
                print(f"  - {file}")
            print("PowerPointレポート生成機能は使用できません。")
            print()

        print("AUS報告書作成アプリケーション（SQLite版）")
        print("=" * 40)
        print(f"ローカルアクセス: http://{Config.HOST}:{Config.PORT}")
        print(f"ネットワークアクセス: http://[このPCのIPアドレス]:{Config.PORT}")
        print()
        print("停止: Ctrl+C")

        app.run(host=Config.HOST, port=Config.PORT, debug=Config.DEBUG)

    except KeyboardInterrupt:
        print("\nアプリケーションを停止しました。")
    except Exception as e:
        print(f"エラーが発生しました: {str(e)}")
        logging.critical(f"アプリケーション起動エラー: {str(e)}")
        input("Enterキーを押して終了してください...")
